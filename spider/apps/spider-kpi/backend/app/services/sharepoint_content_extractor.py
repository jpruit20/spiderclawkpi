"""Deep content extractor for SharePoint files.

Pulls actual content (not just metadata) from every analyzable file
type so the AI analyzer has real text to work with:

- **Excel (.xlsx / .xls)** — every sheet's structure: name, dimensions,
  first 5 rows verbatim (header detection material), then up to 80
  representative data rows. Different from the BOM-only extractor —
  this captures any spreadsheet, BOM or not.
- **PDF** — full text via ``pypdf``, capped at ~30k chars (Claude's
  context budget).
- **Word (.docx)** — paragraph + table text via ``python-docx``.
- **PowerPoint (.pptx)** — slide text + speaker notes via ``python-pptx``.
- **Plain text** — raw bytes decoded as utf-8 (lossy).

Cached in ``sharepoint_file_content`` keyed on document_id.
Re-extracts only when ``source_modified_at`` advances. Sets
``extraction_status='ok' | 'failed' | 'unsupported'`` so downstream
consumers can filter.
"""
from __future__ import annotations

import hashlib
import io
import logging
from datetime import datetime, timezone
from typing import Any, Optional

import requests
from sqlalchemy import select
from sqlalchemy.orm import Session

from app.ingestion.connectors.sharepoint import _get_app_token
from app.models import SharepointDocument, SharepointFileContent


logger = logging.getLogger(__name__)
EXTRACTOR_VERSION = "content-v1.0.0"


# Content limits — Claude has plenty of context but we need to keep
# things bounded so storage + analysis cost stays predictable.
MAX_TEXT_CHARS = 30_000
MAX_SHEET_ROWS = 80
MAX_SHEETS = 30
MAX_PDF_PAGES = 80
MAX_DOCX_PARAGRAPHS = 600
MAX_PPTX_SLIDES = 80


# Mime-type → extractor name dispatch
EXTRACTABLE_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xls",
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "docx",  # best-effort; pypdf-equivalent for .doc not in scope
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    "text/plain": "txt",
}

# Extension fallback when mime is octet-stream (common for SharePoint)
EXTENSION_FALLBACK = {
    ".xlsx": "xlsx", ".xls": "xls",
    ".pdf": "pdf",
    ".docx": "docx", ".doc": "docx",
    ".pptx": "pptx", ".ppt": "pptx",
    ".txt": "txt", ".csv": "txt",
}


def _kind_for(doc: SharepointDocument) -> Optional[str]:
    if doc.mime_type and doc.mime_type in EXTRACTABLE_MIMES:
        return EXTRACTABLE_MIMES[doc.mime_type]
    if doc.name:
        for ext, kind in EXTENSION_FALLBACK.items():
            if doc.name.lower().endswith(ext):
                return kind
    return None


def _download_doc_bytes(doc: SharepointDocument, *, timeout: int = 90) -> bytes:
    token = _get_app_token(doc.tenant_id)
    url = (
        f"https://graph.microsoft.com/v1.0/sites/{doc.graph_site_id}"
        f"/drives/{doc.graph_drive_id}/items/{doc.graph_item_id}/content"
    )
    r = requests.get(url, headers={"Authorization": f"Bearer {token}"}, timeout=timeout, allow_redirects=True)
    r.raise_for_status()
    return r.content


# ── Per-format extractors ──────────────────────────────────────────


def _extract_xlsx(data: bytes) -> tuple[str, dict[str, Any]]:
    import openpyxl

    structure: dict[str, Any] = {"format": "xlsx", "sheets": []}
    text_chunks: list[str] = []
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    for ws in list(wb.worksheets)[:MAX_SHEETS]:
        rows: list[list[Any]] = []
        for r in ws.iter_rows(values_only=True):
            rows.append([("" if c is None else str(c)[:200]) for c in r])
            if len(rows) >= MAX_SHEET_ROWS:
                break
        sheet_summary = {
            "name": ws.title,
            "rows_sampled": len(rows),
            "max_columns_seen": max((len(r) for r in rows), default=0),
            "rows": rows,
        }
        structure["sheets"].append(sheet_summary)
        text_chunks.append(f"=== Sheet: {ws.title} ===")
        for row in rows:
            text_chunks.append(" | ".join(c for c in row if c))
    wb.close()
    text = "\n".join(text_chunks)[:MAX_TEXT_CHARS]
    return text, structure


def _extract_xls(data: bytes) -> tuple[str, dict[str, Any]]:
    import xlrd

    structure: dict[str, Any] = {"format": "xls", "sheets": []}
    text_chunks: list[str] = []
    wb = xlrd.open_workbook(file_contents=data)
    for sheet in wb.sheets()[:MAX_SHEETS]:
        rows: list[list[Any]] = []
        for r in range(min(sheet.nrows, MAX_SHEET_ROWS)):
            rows.append([str(v)[:200] for v in sheet.row_values(r)])
        structure["sheets"].append({
            "name": sheet.name,
            "rows_sampled": len(rows),
            "max_columns_seen": max((len(r) for r in rows), default=0),
            "rows": rows,
        })
        text_chunks.append(f"=== Sheet: {sheet.name} ===")
        for row in rows:
            text_chunks.append(" | ".join(c for c in row if c))
    text = "\n".join(text_chunks)[:MAX_TEXT_CHARS]
    return text, structure


def _extract_pdf(data: bytes) -> tuple[str, dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data), strict=False)
    n_pages = len(reader.pages)
    text_chunks: list[str] = []
    pages_extracted = 0
    for i, page in enumerate(reader.pages[:MAX_PDF_PAGES]):
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t.strip():
            text_chunks.append(f"--- Page {i + 1} ---\n{t}")
            pages_extracted += 1
    text = "\n\n".join(text_chunks)[:MAX_TEXT_CHARS]
    structure = {"format": "pdf", "n_pages": n_pages, "pages_extracted": pages_extracted}
    return text, structure


def _extract_docx(data: bytes) -> tuple[str, dict[str, Any]]:
    import docx

    doc = docx.Document(io.BytesIO(data))
    text_chunks: list[str] = []
    for para in doc.paragraphs[:MAX_DOCX_PARAGRAPHS]:
        if para.text.strip():
            text_chunks.append(para.text)
    # Pull table content too (ECRs, tech packs often use tables)
    table_summaries: list[dict[str, Any]] = []
    for ti, table in enumerate(doc.tables[:30]):
        rows: list[list[str]] = []
        for ri, row in enumerate(table.rows):
            cells = [(c.text or "").strip() for c in row.cells]
            rows.append(cells)
            if ri == 0:
                text_chunks.append(f"=== Table {ti + 1} header: {' | '.join(cells)}")
            else:
                text_chunks.append(" | ".join(cells))
            if ri >= 60:
                break
        table_summaries.append({"table_no": ti + 1, "rows_sampled": len(rows), "rows": rows})
    structure = {"format": "docx", "n_paragraphs": len(doc.paragraphs), "n_tables": len(doc.tables), "tables": table_summaries}
    text = "\n".join(text_chunks)[:MAX_TEXT_CHARS]
    return text, structure


def _extract_pptx(data: bytes) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_chunks: list[str] = []
    slides_summary: list[dict[str, Any]] = []
    for i, slide in enumerate(list(prs.slides)[:MAX_PPTX_SLIDES]):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs = "".join(r.text for r in para.runs)
                    if runs.strip():
                        slide_texts.append(runs.strip())
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""
        text_chunks.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
        if notes_text:
            text_chunks.append(f"  [notes] {notes_text}")
        slides_summary.append({"slide_no": i + 1, "lines": len(slide_texts), "has_notes": bool(notes_text)})
    structure = {"format": "pptx", "n_slides": len(prs.slides), "slides": slides_summary}
    text = "\n\n".join(text_chunks)[:MAX_TEXT_CHARS]
    return text, structure


def _extract_txt(data: bytes) -> tuple[str, dict[str, Any]]:
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = data.decode("latin-1", errors="replace")
    return text[:MAX_TEXT_CHARS], {"format": "txt", "raw_bytes": len(data)}


_DISPATCH = {
    "xlsx": _extract_xlsx,
    "xls": _extract_xls,
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "txt": _extract_txt,
}


# ── Public API ──────────────────────────────────────────────────────


def extract_content_for_document(db: Session, doc: SharepointDocument, *, force: bool = False) -> dict[str, Any]:
    """Fetch + parse a single document. Returns the persisted row's
    summary status. Re-uses cached row if source_modified_at hasn't
    advanced (cheap idempotency for daily sweeps)."""
    kind = _kind_for(doc)
    if kind is None:
        # Mark as unsupported so we don't keep re-checking
        return _persist_status(db, doc, status="unsupported", error="no extractor for mime/extension")

    existing = db.execute(
        select(SharepointFileContent).where(SharepointFileContent.document_id == doc.id)
    ).scalar_one_or_none()

    if (
        existing is not None
        and not force
        and existing.extractor_version == EXTRACTOR_VERSION
        and existing.source_modified_at is not None
        and doc.modified_at_remote is not None
        and existing.source_modified_at >= doc.modified_at_remote
        and existing.extraction_status == "ok"
    ):
        return {"status": "cached", "id": existing.id}

    try:
        data = _download_doc_bytes(doc)
    except Exception as exc:
        return _persist_status(db, doc, status="failed", error=f"download: {exc}"[:1024])

    sha = hashlib.sha256(data).hexdigest()
    extractor = _DISPATCH[kind]
    try:
        text, structure = extractor(data)
    except Exception as exc:
        logger.exception("content extraction failed for doc %s (%s)", doc.id, doc.name)
        return _persist_status(db, doc, status="failed", error=f"parse({kind}): {exc}"[:1024], byte_size=len(data), sha=sha)

    return _persist_content(
        db,
        doc,
        text=text,
        structure=structure,
        sha=sha,
        byte_size=len(data),
    )


def _persist_status(
    db: Session,
    doc: SharepointDocument,
    *,
    status: str,
    error: Optional[str] = None,
    byte_size: Optional[int] = None,
    sha: Optional[str] = None,
) -> dict[str, Any]:
    existing = db.execute(
        select(SharepointFileContent).where(SharepointFileContent.document_id == doc.id)
    ).scalar_one_or_none()
    now = datetime.now(timezone.utc)
    if existing is None:
        existing = SharepointFileContent(
            document_id=doc.id,
            extraction_status=status,
            extraction_error=error,
            source_modified_at=doc.modified_at_remote,
            byte_size=byte_size,
            content_sha256=sha,
            extractor_version=EXTRACTOR_VERSION,
            extracted_at=now,
        )
        db.add(existing)
    else:
        existing.extraction_status = status
        existing.extraction_error = error
        existing.source_modified_at = doc.modified_at_remote
        existing.byte_size = byte_size if byte_size is not None else existing.byte_size
        existing.content_sha256 = sha if sha is not None else existing.content_sha256
        existing.extractor_version = EXTRACTOR_VERSION
        existing.extracted_at = now
    db.commit()
    return {"status": status, "id": existing.id, "error": error}


def _persist_content(
    db: Session,
    doc: SharepointDocument,
    *,
    text: str,
    structure: dict[str, Any],
    sha: str,
    byte_size: int,
) -> dict[str, Any]:
    existing = db.execute(
        select(SharepointFileContent).where(SharepointFileContent.document_id == doc.id)
    ).scalar_one_or_none()
    now = datetime.now(timezone.utc)
    if existing is None:
        existing = SharepointFileContent(
            document_id=doc.id,
            text_content=text,
            structure_json=structure,
            content_sha256=sha,
            byte_size=byte_size,
            source_modified_at=doc.modified_at_remote,
            extractor_version=EXTRACTOR_VERSION,
            extraction_status="ok",
            extracted_at=now,
        )
        db.add(existing)
    else:
        existing.text_content = text
        existing.structure_json = structure
        existing.content_sha256 = sha
        existing.byte_size = byte_size
        existing.source_modified_at = doc.modified_at_remote
        existing.extractor_version = EXTRACTOR_VERSION
        existing.extraction_status = "ok"
        existing.extraction_error = None
        existing.extracted_at = now
    db.commit()
    return {"status": "ok", "id": existing.id, "text_chars": len(text)}


# Categories that are worth deep-extracting (skip pure images, video,
# CAD binary blobs that have no meaningful text).
ANALYSIS_WORTHY_SEMANTICS = {
    "bom", "cbom", "price_list", "tech_pack", "vendor_doc", "ecr",
    "design_doc", "drawing", "manual", "test_report", "qa_doc",
    "firmware_doc", "packaging", "label_or_cert",
    "spreadsheet", "pdf", "word_doc", "presentation",
}


def extract_content_for_corpus(
    db: Session,
    *,
    spider_product: Optional[str] = None,
    archive_status: str = "active",
    limit: Optional[int] = None,
    force: bool = False,
) -> dict[str, int]:
    """Walk every analysis-worthy active doc, extracting content for any
    that don't have a fresh ``sharepoint_file_content`` row."""
    q = (
        select(SharepointDocument)
        .where(
            SharepointDocument.is_folder == False,  # noqa: E712
            SharepointDocument.archive_status == archive_status,
            SharepointDocument.semantic_type.in_(ANALYSIS_WORTHY_SEMANTICS),
        )
        .order_by(SharepointDocument.modified_at_remote.desc().nulls_last())
    )
    if spider_product:
        q = q.where(SharepointDocument.spider_product == spider_product)
    if limit:
        q = q.limit(limit)

    docs = db.execute(q).scalars().all()
    counts = {"seen": 0, "extracted": 0, "cached": 0, "failed": 0, "unsupported": 0}
    for doc in docs:
        counts["seen"] += 1
        result = extract_content_for_document(db, doc, force=force)
        s = result.get("status")
        if s == "cached":
            counts["cached"] += 1
        elif s == "ok":
            counts["extracted"] += 1
        elif s == "unsupported":
            counts["unsupported"] += 1
        else:
            counts["failed"] += 1
    return counts
